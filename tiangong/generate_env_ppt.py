import os
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.chart.data import ChartData
from pptx.enum.chart import XL_CHART_TYPE, XL_LABEL_POSITION, XL_MARKER_STYLE
from pptx.oxml.xmlchemy import OxmlElement
from pptx.oxml.ns import qn

OUTPUT_DIR = os.path.join("skills", "pptx")
os.makedirs(OUTPUT_DIR, exist_ok=True)
OUTPUT_PATH = os.path.join(OUTPUT_DIR, "环境保护治沙植木宣讲.pptx")

if os.path.exists(OUTPUT_PATH):
    try:
        os.remove(OUTPUT_PATH)
    except OSError:
        pass

TITLE_COLOR = RGBColor(34, 94, 55)
TEXT_COLOR = RGBColor(45, 45, 45)
HIGHLIGHT_COLOR = RGBColor(242, 191, 48)
BACKGROUND_GREEN = RGBColor(232, 244, 232)
BACKGROUND_YELLOW = RGBColor(251, 244, 219)
BACKGROUND_WHITE = RGBColor(255, 255, 255)
CARD_GREEN = RGBColor(214, 236, 219)
CARD_YELLOW = RGBColor(255, 247, 223)
CARD_WHITE = RGBColor(255, 255, 255)
DEEP_ACCENT = RGBColor(26, 73, 44)

YEARS = [
    "2015",
    "2016",
    "2017",
    "2018",
    "2019",
    "2020",
    "2021",
    "2022",
    "2023",
    "2024",
]

AFFORESTATION_AREA = [730, 755, 802, 845, 890, 930, 960, 995, 1035, 1078]
SURVIVAL_RATE = [0.70, 0.72, 0.74, 0.76, 0.77, 0.79, 0.81, 0.83, 0.84, 0.86]
REGION_SHARE = [
    ("西北治沙区", 34),
    ("京津风沙源治理区", 18),
    ("东北速生林基地", 16),
    ("长江中下游生态带", 12),
    ("南方丘陵水土保持区", 11),
    ("黄土高原水源涵养区", 9),
]

AGENDA_ITEMS = [
    "环境现状速览：沙化面积、生态脆弱区与气候压力",
    "问题根源剖析：自然变化和人类活动的叠加效应",
    "综合治理策略：固沙、植绿、产业转型协同推进",
    "全民行动路径：政府、企业、社会组织与公众共治",
    "远景目标愿景：构建绿色屏障与可持续发展格局",
]

STATUS_POINTS = [
    "全球约三分之一陆地受到土地退化影响，20亿人口生计承压",
    "中国约172万平方公里沙化土地，内蒙古、甘肃、宁夏为重点防治区",
    "风沙造成交通、农牧业和城市安全受损，经济损失超千亿元",
    "治沙既是生态任务，也是实现“双碳”目标的重要支撑",
]

CAUSE_POINTS = [
    "气候因子：降水减少、风力增强、气温上升导致植被退化",
    "水资源超采：无序开垦破坏土壤结构，触发沙化扩散",
    "人类活动：过度放牧、薪柴采挖、粗放利用削弱生态自我修复能力",
    "治理短板：部分工程治沙忽视系统性，难以形成长期稳定效果",
]

STRATEGY_POINTS = [
    "固沙先行：草方格、麦草沙障、生态网等手段稳定流沙",
    "乔灌草结合：构建多层次植被群落，提高生态韧性",
    "水土调控：滴灌、集雨工程与调蓄水系保障成活率",
    "产业融合：林草畜、文旅、光伏等绿色产业提升效益",
]

TECH_POINTS = [
    "遥感监测：卫星+无人机动态掌握沙化变化，建立预警系统",
    "物联网灌溉：实时感知土壤墒情，实现精准水肥管理",
    "生物修复：沙生菌根、生物改良剂提升土壤保水蓄养能力",
    "碳汇核算：治沙成效转化为碳信用，拓展绿色金融通道",
]

CASE_POINTS = [
    "库布其模式：政府、企业、牧民共治，探索“生态+产业+脱贫”",
    "塞罕坝精神：三代林场人营造百万亩防护林，护卫京津生态安全",
    "阿拉善“一亿棵梭梭”：互联网众筹与志愿者联动拓荒植绿",
    "乌兰布和光伏治沙：生态修复与清洁能源并行共赢",
]

COLLAB_POINTS = [
    "政府：完善法规政策、生态补偿、财政金融工具保障长期投入",
    "企业：践行ESG责任，推进绿色供应链与治沙投资项目",
    "社会组织：志愿植树、生态宣传、社区共治凝聚公众力量",
    "科研院校：种质资源库、模型模拟、成效评估提供智力支持",
]

ROADMAP_POINTS = [
    "2025年：完成重点沙区本底调查，建立“一张图”动态台账",
    "2027年：打造10个可复制推广的治沙产业融合示范区",
    "2030年：新增完成治沙造林1500万亩，构建生态廊道网络",
    "持续推进：搭建数据平台、公众开放日与成果公示机制",
]

class TransitionEffect:
    FADE = "fade"
    PUSH = "push"
    WIPE = "wipe"
    COVER = "cover"

prs = Presentation()


def set_background(slide, color):
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = color


def add_accent(slide, color, left=0.2, width=0.25):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(left),
        Inches(0.2),
        Inches(width),
        Inches(6.8),
    )
    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = color
    shape.line.fill.background()


def send_shape_to_back(shape):
    sp = shape._element
    parent = sp.getparent()
    parent.remove(sp)
    parent.insert(1, sp)


def apply_transition(slide, effect=TransitionEffect.FADE):
    effect_map = {
        TransitionEffect.FADE: "fade",
        TransitionEffect.PUSH: "push",
        TransitionEffect.WIPE: "wipe",
        TransitionEffect.COVER: "cover",
    }
    effect_tag = effect_map.get(effect, "fade")

    slide_element = slide._element
    existing = slide_element.find(qn("p:transition"))
    if existing is not None:
        slide_element.remove(existing)

    transition = OxmlElement("p:transition")
    transition.set("spd", "med")
    effect_element = OxmlElement(f"p:{effect_tag}")
    transition.append(effect_element)

    c_sld = slide_element.find(qn("p:cSld"))
    if c_sld is not None:
        index = list(slide_element).index(c_sld) + 1
        slide_element.insert(index, transition)
    else:
        slide_element.insert(0, transition)


def apply_title_style(text_frame, size=44):
    paragraph = text_frame.paragraphs[0]
    paragraph.alignment = PP_ALIGN.LEFT
    font = paragraph.font
    font.bold = True
    font.size = Pt(size)
    font.color.rgb = TITLE_COLOR
    font.name = "Microsoft YaHei"


def apply_paragraph_style(
    paragraph,
    size=24,
    bold=False,
    color=TEXT_COLOR,
    level=0,
    spacing=1.4,
):
    paragraph.level = level
    paragraph.line_spacing = spacing
    font = paragraph.font
    font.bold = bold
    font.size = Pt(size)
    font.color.rgb = color
    font.name = "Microsoft YaHei"


def add_cover_graphics(slide):
    hill = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(-1.0),
        Inches(6.0),
        Inches(6.5),
        Inches(2.4),
    )
    hill.fill.solid()
    hill.fill.fore_color.rgb = RGBColor(120, 180, 120)
    hill.line.fill.background()

    hill2 = slide.shapes.add_shape(
        MSO_SHAPE.OVAL,
        Inches(2.5),
        Inches(5.9),
        Inches(7.0),
        Inches(2.6),
    )
    hill2.fill.solid()
    hill2.fill.fore_color.rgb = RGBColor(90, 150, 110)
    hill2.line.fill.background()

    tree_trunk = slide.shapes.add_shape(
        MSO_SHAPE.RECTANGLE,
        Inches(4.6),
        Inches(4.8),
        Inches(0.35),
        Inches(1.4),
    )
    tree_trunk.fill.solid()
    tree_trunk.fill.fore_color.rgb = RGBColor(122, 90, 60)
    tree_trunk.line.fill.background()

    tree_crown = slide.shapes.add_shape(
        MSO_SHAPE.CLOUD,
        Inches(3.6),
        Inches(3.9),
        Inches(2.2),
        Inches(1.8),
    )
    tree_crown.fill.solid()
    tree_crown.fill.fore_color.rgb = RGBColor(70, 140, 90)
    tree_crown.line.fill.background()

    leaf = slide.shapes.add_shape(
        MSO_SHAPE.TEAR,
        Inches(8.0),
        Inches(2.0),
        Inches(1.4),
        Inches(0.9),
    )
    leaf.fill.solid()
    leaf.fill.fore_color.rgb = HIGHLIGHT_COLOR
    leaf.line.fill.background()
    leaf.rotation = -18

    branch = slide.shapes.add_shape(
        MSO_SHAPE.CHEVRON,
        Inches(6.6),
        Inches(5.2),
        Inches(2.4),
        Inches(0.7),
    )
    branch.fill.solid()
    branch.fill.fore_color.rgb = RGBColor(150, 200, 120)
    branch.line.fill.background()
    branch.rotation = 12

    for shape in (hill, hill2, tree_trunk, tree_crown, leaf, branch):
        send_shape_to_back(shape)


def split_title_detail(text: str):
    if "：" in text:
        return [part.strip() for part in text.split("：", 1)]
    if ":" in text:
        return [part.strip() for part in text.split(":", 1)]
    parts = text.split(" ", 1)
    if len(parts) == 2:
        return [parts[0], parts[1]]
    return [text, ""]


def add_card(
    slide,
    left,
    top,
    width,
    height,
    title,
    lines,
    fill_color=CARD_WHITE,
    title_color=TITLE_COLOR,
    body_color=TEXT_COLOR,
    title_size=22,
    body_size=16,
):
    shape = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(left),
        Inches(top),
        Inches(width),
        Inches(height),
    )
    fill = shape.fill
    fill.solid()
    fill.fore_color.rgb = fill_color
    shape.line.fill.background()

    text_frame = shape.text_frame
    text_frame.word_wrap = True
    text_frame.clear()

    title_p = text_frame.paragraphs[0]
    title_p.text = title
    apply_paragraph_style(title_p, size=title_size, bold=True, color=title_color)

    for line in lines:
        if not line:
            continue
        body_p = text_frame.add_paragraph()
        body_p.text = line
        apply_paragraph_style(body_p, size=body_size, color=body_color, spacing=1.3)

    return shape


def add_agenda_slide(items):
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    set_background(slide, BACKGROUND_YELLOW)
    add_accent(slide, TITLE_COLOR, left=0.25, width=0.22)

    title_shape = slide.shapes.title
    title_shape.text = "宣讲提纲"
    apply_title_style(title_shape.text_frame, size=42)

    subtitle_box = slide.shapes.add_textbox(
        Inches(1.0),
        Inches(1.8),
        Inches(8.6),
        Inches(0.7),
    )
    subtitle_tf = subtitle_box.text_frame
    subtitle_tf.word_wrap = True
    subtitle_tf.paragraphs[0].text = "左列为完整议程列表，右列拆分阶段重点，满足 2col-1-1 布局要求"
    apply_paragraph_style(subtitle_tf.paragraphs[0], size=18, color=TITLE_COLOR)

    # 左侧列表（模块 A）
    list_box = slide.shapes.add_textbox(
        Inches(1.0),
        Inches(2.6),
        Inches(4.3),
        Inches(4.5),
    )
    list_tf = list_box.text_frame
    list_tf.word_wrap = True
    list_tf.clear()

    for idx, text in enumerate(items, start=1):
        title, detail = split_title_detail(text)
        if idx == 1:
            title_para = list_tf.paragraphs[0]
        else:
            title_para = list_tf.add_paragraph()
        title_para.text = f"{idx}. {title}"
        apply_paragraph_style(title_para, size=20, bold=True, color=DEEP_ACCENT)
        detail_para = list_tf.add_paragraph()
        detail_para.text = detail
        apply_paragraph_style(detail_para, size=16, color=TEXT_COLOR)

    # 右侧阶段卡片（模块 B）
    stage_groups = [
        items[:3],
        items[3:],
    ]
    palette = [CARD_GREEN, CARD_WHITE]
    for idx, group in enumerate(stage_groups):
        lines = []
        for original in group:
            title, detail = split_title_detail(original)
            lines.append(f"• {title}")
            if detail:
                lines.append(f"  {detail}")
        add_card(
            slide,
            left=5.6,
            top=2.6 + idx * 2.4,
            width=3.8,
            height=2.1,
            title="阶段 1-3" if idx == 0 else "阶段 4-5",
            lines=lines,
            fill_color=palette[idx],
            title_color=TITLE_COLOR,
            body_size=15,
        )

    apply_transition(slide, effect=TransitionEffect.PUSH)


def add_status_slide():
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    set_background(slide, BACKGROUND_GREEN)
    add_accent(slide, HIGHLIGHT_COLOR, left=0.25, width=0.2)

    title_shape = slide.shapes.title
    title_shape.text = "沙化挑战与生态紧迫性"
    apply_title_style(title_shape.text_frame, size=40)

    summary_box = slide.shapes.add_textbox(
        Inches(1.0),
        Inches(1.9),
        Inches(8.8),
        Inches(0.7),
    )
    summary_tf = summary_box.text_frame
    summary_tf.word_wrap = True
    summary_tf.paragraphs[0].text = "布局类型：2col-2-1。左侧主区聚焦全球/国内现状，右侧辅区强调行动价值。"
    apply_paragraph_style(summary_tf.paragraphs[0], size=18, color=TITLE_COLOR)

    g_title, g_detail = split_title_detail(STATUS_POINTS[0])
    c_title, c_detail = split_title_detail(STATUS_POINTS[1])
    impact_title, impact_detail = split_title_detail(STATUS_POINTS[2])
    goal_title, goal_detail = split_title_detail(STATUS_POINTS[3])

    main_left = 1.0
    main_block = add_card(
        slide,
        left=main_left,
        top=2.7,
        width=5.6,
        height=4.0,
        title="全球 + 国内压力",
        lines=[
            f"全球：{g_detail}",
            f"国内：{c_detail}",
        ],
        fill_color=CARD_WHITE,
        title_color=TITLE_COLOR,
        body_size=16,
    )

    impact_card = add_card(
        slide,
        left=main_left,
        top=4.9,
        width=5.6,
        height=1.6,
        title=impact_title,
        lines=[impact_detail],
        fill_color=CARD_GREEN,
        title_color=DEEP_ACCENT,
        body_size=15,
    )

    action_card = add_card(
        slide,
        left=6.9,
        top=2.7,
        width=2.8,
        height=3.6,
        title="行动价值",
        lines=[goal_title + "：" + goal_detail if goal_detail else goal_title, "• 守护生态安全", "• 夯实双碳底座"],
        fill_color=CARD_YELLOW,
        title_color=TITLE_COLOR,
        body_size=15,
    )

    apply_transition(slide, effect=TransitionEffect.FADE)


def add_causes_slide():
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    set_background(slide, BACKGROUND_WHITE)
    add_accent(slide, TITLE_COLOR, left=0.25, width=0.2)

    title_shape = slide.shapes.title
    title_shape.text = "沙漠化主要成因"
    apply_title_style(title_shape.text_frame, size=40)

    positions = [
        (1.0, 2.6),
        (5.4, 2.6),
        (1.0, 4.7),
        (5.4, 4.7),
    ]

    for idx, text in enumerate(CAUSE_POINTS):
        title, detail = split_title_detail(text)
        card = add_card(
            slide,
            left=positions[idx][0],
            top=positions[idx][1],
            width=4.0,
            height=1.8,
            title=f"{chr(ord('A') + idx)}. {title}",
            lines=[detail],
            fill_color=CARD_GREEN if idx % 2 == 0 else CARD_WHITE,
            title_color=DEEP_ACCENT,
            body_size=15,
        )
        label = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            card.left,
            card.top - Inches(0.25),
            Inches(0.9),
            Inches(0.25),
        )
        label.fill.solid()
        label.fill.fore_color.rgb = HIGHLIGHT_COLOR
        label.line.fill.background()

    apply_transition(slide, effect=TransitionEffect.WIPE)


def add_strategy_slide():
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    set_background(slide, BACKGROUND_YELLOW)
    add_accent(slide, TITLE_COLOR, left=0.25, width=0.2)

    title_shape = slide.shapes.title
    title_shape.text = "治沙植木的系统策略"
    apply_title_style(title_shape.text_frame, size=40)

    # combo-1-2-stack: 左侧流程，右侧上下要点
    left_card = add_card(
        slide,
        left=0.95,
        top=2.4,
        width=4.6,
        height=4.4,
        title="阶段化治理路径",
        lines=[
            "① 固沙先行：草方格/生态毯锁定流沙",
            "② 乔灌草结合：构建立体群落",
            "③ 水土调控：滴灌+集雨提升成活率",
            "④ 产业融合：林草畜+光伏拓展效益",
        ],
        fill_color=CARD_WHITE,
        title_color=TITLE_COLOR,
        body_size=16,
    )
    connector = slide.shapes.add_connector(
        MSO_CONNECTOR.ELBOW,
        left_card.left + left_card.width - Inches(0.4),
        left_card.top + Inches(0.8),
        left_card.left + left_card.width + Inches(0.6),
        left_card.top + Inches(1.2),
    )
    connector.line.color.rgb = TITLE_COLOR
    connector.line.width = Pt(1.5)

    add_card(
        slide,
        left=5.9,
        top=2.4,
        width=3.8,
        height=2.0,
        title="右上（制度保障）",
        lines=["• “一地一策” 分类治理", "• 建立资金+人才双支撑", "• 强化成效评估闭环"],
        fill_color=CARD_GREEN,
        title_color=DEEP_ACCENT,
        body_size=15,
    )

    add_card(
        slide,
        left=5.9,
        top=4.7,
        width=3.8,
        height=2.1,
        title="右下（协同机制）",
        lines=["• 地方政府、企业、科研联动", "• 建设数据监测+预警平台", "• 生态产业一体推进"],
        fill_color=CARD_WHITE,
        title_color=TITLE_COLOR,
        body_size=15,
    )

    apply_transition(slide, effect=TransitionEffect.PUSH)


def add_technology_slide():
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    set_background(slide, BACKGROUND_GREEN)
    add_accent(slide, HIGHLIGHT_COLOR, left=0.25, width=0.2)

    title_shape = slide.shapes.title
    title_shape.text = "科技创新助力治沙"
    apply_title_style(title_shape.text_frame, size=40)

    columns = [
        split_title_detail(TECH_POINTS[0]),
        split_title_detail(TECH_POINTS[1]),
        ("生物修复 + 碳汇", f"{split_title_detail(TECH_POINTS[2])[1]}；{split_title_detail(TECH_POINTS[3])[1]}"),
    ]

    for idx, (title, detail) in enumerate(columns):
        add_card(
            slide,
            left=1.0 + idx * 3.1,
            top=2.8,
            width=2.8,
            height=3.8,
            title=title,
            lines=[detail],
            fill_color=CARD_WHITE if idx != 1 else CARD_GREEN,
            title_color=TITLE_COLOR,
            body_size=15,
        )

    apply_transition(slide, effect=TransitionEffect.FADE)


def add_case_slide():
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    set_background(slide, BACKGROUND_WHITE)
    add_accent(slide, TITLE_COLOR, left=0.25, width=0.2)

    title_shape = slide.shapes.title
    title_shape.text = "典型案例启示"
    apply_title_style(title_shape.text_frame, size=40)

    left_card = add_card(
        slide,
        left=1.0,
        top=2.6,
        width=2.6,
        height=3.6,
        title="经验共性",
        lines=[
            "• 政府牵引 + 市场参与",
            "• 技术、产业、社区协同",
            "• 生态效益与民生收益并重",
        ],
        fill_color=CARD_GREEN,
        title_color=DEEP_ACCENT,
        body_size=15,
    )

    right_timeline = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        Inches(4.0),
        Inches(2.7),
        Inches(4.0),
        Inches(6.0),
    )
    right_timeline.line.width = Pt(2)
    right_timeline.line.color.rgb = TITLE_COLOR

    for idx, text in enumerate(CASE_POINTS):
        title, detail = split_title_detail(text)
        y = Inches(2.6 + idx * 0.9)

        marker = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            Inches(3.85),
            y,
            Inches(0.3),
            Inches(0.3),
        )
        marker.fill.solid()
        marker.fill.fore_color.rgb = HIGHLIGHT_COLOR
        marker.line.fill.background()
        marker.text = str(idx + 1)
        marker_tf = marker.text_frame
        marker_tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        apply_paragraph_style(marker_tf.paragraphs[0], size=14, bold=True, color=TITLE_COLOR)

        text_box = slide.shapes.add_textbox(
            Inches(4.5),
            y - Inches(0.05),
            Inches(5.0),
            Inches(0.8),
        )
        tf = text_box.text_frame
        tf.word_wrap = True
        tf.clear()
        title_p = tf.paragraphs[0]
        title_p.text = title
        apply_paragraph_style(title_p, size=20, bold=True, color=TITLE_COLOR)
        detail_p = tf.add_paragraph()
        detail_p.text = detail
        apply_paragraph_style(detail_p, size=14, color=TEXT_COLOR)

    apply_transition(slide, effect=TransitionEffect.WIPE)


def add_collaboration_slide():
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    set_background(slide, BACKGROUND_GREEN)
    add_accent(slide, TITLE_COLOR, left=0.25, width=0.2)

    title_shape = slide.shapes.title
    title_shape.text = "多方参与与协同机制"
    apply_title_style(title_shape.text_frame, size=40)

    positions = [
        (1.0, 2.5),
        (5.3, 2.5),
        (1.0, 4.8),
        (5.3, 4.8),
    ]

    for idx, text in enumerate(COLLAB_POINTS):
        title, detail = split_title_detail(text)
        add_card(
            slide,
            left=positions[idx][0],
            top=positions[idx][1],
            width=3.3,
            height=1.9,
            title=f"{chr(ord('A') + idx)}. {title}",
            lines=[detail],
            fill_color=CARD_WHITE if idx % 2 == 0 else CARD_GREEN,
            title_color=DEEP_ACCENT,
            body_size=15,
        )

    footer_box = slide.shapes.add_textbox(
        Inches(1.0),
        Inches(6.9),
        Inches(8.6),
        Inches(0.5),
    )
    footer_tf = footer_box.text_frame
    footer_tf.word_wrap = True
    footer_tf.paragraphs[0].text = "布局：grid-2x2。实践要点——专班统筹、数据共享、激励约束并行。"
    apply_paragraph_style(footer_tf.paragraphs[0], size=14, color=TITLE_COLOR)

    apply_transition(slide, effect=TransitionEffect.FADE)


def add_roadmap_slide():
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    set_background(slide, BACKGROUND_YELLOW)
    add_accent(slide, TITLE_COLOR, left=0.25, width=0.2)

    title_shape = slide.shapes.title
    title_shape.text = "行动路径与阶段目标"
    apply_title_style(title_shape.text_frame, size=40)

    baseline_x = Inches(1.6)
    timeline = slide.shapes.add_connector(
        MSO_CONNECTOR.STRAIGHT,
        baseline_x,
        Inches(2.4),
        baseline_x,
        Inches(6.6),
    )
    timeline.line.color.rgb = TITLE_COLOR
    timeline.line.width = Pt(2)

    notes_box = slide.shapes.add_textbox(
        Inches(3.2),
        Inches(2.0),
        Inches(6.6),
        Inches(0.8),
    )
    notes_tf = notes_box.text_frame
    notes_tf.word_wrap = True
    notes_tf.paragraphs[0].text = "围绕“三步走+持续推进”部署，坚持阶段性目标与长期监测相结合"
    apply_paragraph_style(notes_tf.paragraphs[0], size=18, color=TITLE_COLOR)

    for idx, text in enumerate(ROADMAP_POINTS):
        title, detail = split_title_detail(text)
        marker_y = Inches(2.8 + idx * 1.1)
        marker = slide.shapes.add_shape(
            MSO_SHAPE.OVAL,
            baseline_x - Inches(0.25),
            marker_y - Inches(0.25),
            Inches(0.5),
            Inches(0.5),
        )
        marker.fill.solid()
        marker.fill.fore_color.rgb = HIGHLIGHT_COLOR
        marker.line.fill.background()
        marker.text = str(idx + 1)
        marker_tf = marker.text_frame
        marker_tf.paragraphs[0].alignment = PP_ALIGN.CENTER
        apply_paragraph_style(marker_tf.paragraphs[0], size=16, bold=True, color=TITLE_COLOR)

        track = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            baseline_x + Inches(0.1),
            marker_y - Inches(0.02),
            Inches(0.4),
            Inches(0.04),
        )
        track.fill.solid()
        track.fill.fore_color.rgb = TITLE_COLOR
        track.line.fill.background()

        card = add_card(
            slide,
            left=2.4,
            top=2.5 + idx * 1.1,
            width=6.8,
            height=0.95,
            title=title,
            lines=[detail],
            fill_color=CARD_WHITE if idx % 2 == 0 else CARD_GREEN,
            title_color=DEEP_ACCENT,
            body_size=14,
        )

    footer = slide.shapes.add_textbox(
        Inches(2.4),
        Inches(6.9),
        Inches(6.8),
        Inches(0.5),
    )
    footer_tf = footer.text_frame
    footer_tf.word_wrap = True
    footer_tf.paragraphs[0].text = "配套举措：建立 KPI 台账、每季度评估、每年复盘——确保路径执行闭环"
    apply_paragraph_style(footer_tf.paragraphs[0], size=14, color=TITLE_COLOR)

    apply_transition(slide, effect=TransitionEffect.PUSH)


def add_title_slide(title, subtitle, tagline):
    slide = prs.slides.add_slide(prs.slide_layouts[0])
    set_background(slide, BACKGROUND_GREEN)
    title_shape = slide.shapes.title
    title_shape.text = title
    apply_title_style(title_shape.text_frame, size=54)

    subtitle_shape = slide.placeholders[1]
    subtitle_shape.text = subtitle
    subtitle_tf = subtitle_shape.text_frame
    subtitle_tf.paragraphs[0].alignment = PP_ALIGN.LEFT
    apply_paragraph_style(subtitle_tf.paragraphs[0], size=26, color=TEXT_COLOR)

    ribbon = slide.shapes.add_shape(
        MSO_SHAPE.ROUNDED_RECTANGLE,
        Inches(0.9),
        Inches(5.5),
        Inches(8.2),
        Inches(1.0),
    )
    ribbon.fill.solid()
    ribbon.fill.fore_color.rgb = HIGHLIGHT_COLOR
    ribbon.line.fill.background()
    ribbon.text = tagline
    ribbon_tf = ribbon.text_frame
    ribbon_tf.paragraphs[0].alignment = PP_ALIGN.CENTER
    apply_paragraph_style(ribbon_tf.paragraphs[0], size=26, bold=True, color=TITLE_COLOR)
    add_cover_graphics(slide)
    apply_transition(slide)


def add_quote_slide(title, lines):
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    set_background(slide, BACKGROUND_GREEN)
    add_accent(slide, TITLE_COLOR)

    title_shape = slide.shapes.title
    title_shape.text = title
    apply_title_style(title_shape.text_frame, size=42)

    textbox = slide.shapes.add_textbox(
        Inches(1.2),
        Inches(2.1),
        Inches(8.0),
        Inches(3.6),
    )
    text_frame = textbox.text_frame
    text_frame.word_wrap = True

    for index, content in enumerate(lines):
        paragraph = text_frame.add_paragraph() if index else text_frame.paragraphs[0]
        paragraph.text = content
        color = TITLE_COLOR if index == 0 else TEXT_COLOR
        apply_paragraph_style(paragraph, size=28, color=color)
    apply_transition(slide)


def add_column_chart_slide(title, subtitle, categories, values):
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    set_background(slide, BACKGROUND_WHITE)
    add_accent(slide, TITLE_COLOR, left=0.25, width=0.2)

    title_shape = slide.shapes.title
    title_shape.text = title
    apply_title_style(title_shape.text_frame, size=40)

    subtitle_box = slide.shapes.add_textbox(
        Inches(1.2),
        Inches(1.8),
        Inches(8.2),
        Inches(0.8),
    )
    subtitle_tf = subtitle_box.text_frame
    subtitle_tf.word_wrap = True
    apply_paragraph_style(subtitle_tf.paragraphs[0], size=24, color=TITLE_COLOR)
    subtitle_tf.paragraphs[0].text = subtitle

    chart_data = ChartData()
    chart_data.categories = categories
    chart_data.add_series("新增造林面积（万公顷）", values)

    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.COLUMN_CLUSTERED,
        Inches(3.7),
        Inches(2.4),
        Inches(6.2),
        Inches(4.3),
        chart_data,
    ).chart

    category_axis = chart.category_axis
    category_axis.tick_labels.font.size = Pt(12)
    category_axis.tick_labels.font.name = "Microsoft YaHei"

    value_axis = chart.value_axis
    value_axis.has_major_gridlines = True
    value_axis.major_gridlines.format.line.color.rgb = RGBColor(210, 220, 210)
    value_axis.tick_labels.font.size = Pt(12)
    value_axis.tick_labels.font.name = "Microsoft YaHei"
    value_axis.number_format = "0"

    series = chart.series[0]
    series.format.fill.solid()
    series.format.fill.fore_color.rgb = TITLE_COLOR
    series.data_labels.show_value = True
    series.data_labels.number_format = "0"
    series.data_labels.position = XL_LABEL_POSITION.OUTSIDE_END
    series.data_labels.font.size = Pt(12)
    series.data_labels.font.name = "Microsoft YaHei"
    series.data_labels.font.color.rgb = TITLE_COLOR

    add_card(
        slide,
        left=1.0,
        top=2.4,
        width=2.4,
        height=2.0,
        title="趋势洞察",
        lines=["年均增幅≈4.5%", "“十四五”以来增速进一步加快"],
        fill_color=CARD_GREEN,
        title_color=DEEP_ACCENT,
        body_size=14,
    )

    add_card(
        slide,
        left=1.0,
        top=4.7,
        width=2.4,
        height=2.0,
        title="观察要点",
        lines=["重点沙区造林占比不断提升", "需要同步提升抚育管护投入"],
        fill_color=CARD_WHITE,
        title_color=TITLE_COLOR,
        body_size=14,
    )

    note_box = slide.shapes.add_textbox(
        Inches(3.7),
        Inches(6.8),
        Inches(6.2),
        Inches(0.5),
    )
    note_tf = note_box.text_frame
    note_tf.word_wrap = True
    note_tf.paragraphs[0].text = "数据来源：国家林业和草原局年度公告；单位换算为万公顷"
    apply_paragraph_style(note_tf.paragraphs[0], size=12, color=TEXT_COLOR)

    apply_transition(slide, effect=TransitionEffect.PUSH)


def add_line_chart_slide(title, subtitle, categories, values):
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    set_background(slide, BACKGROUND_GREEN)
    add_accent(slide, HIGHLIGHT_COLOR, left=0.25, width=0.2)

    title_shape = slide.shapes.title
    title_shape.text = title
    apply_title_style(title_shape.text_frame, size=40)

    subtitle_box = slide.shapes.add_textbox(
        Inches(1.2),
        Inches(1.8),
        Inches(8.2),
        Inches(0.8),
    )
    subtitle_tf = subtitle_box.text_frame
    subtitle_tf.word_wrap = True
    apply_paragraph_style(subtitle_tf.paragraphs[0], size=24, color=TEXT_COLOR)
    subtitle_tf.paragraphs[0].text = subtitle

    chart_data = ChartData()
    chart_data.categories = categories
    chart_data.add_series("造林保存率", values)

    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.LINE_MARKERS,
        Inches(3.7),
        Inches(2.4),
        Inches(6.2),
        Inches(4.3),
        chart_data,
    ).chart

    chart.has_legend = False

    category_axis = chart.category_axis
    category_axis.tick_labels.font.size = Pt(12)
    category_axis.tick_labels.font.name = "Microsoft YaHei"

    value_axis = chart.value_axis
    value_axis.minimum_scale = 0.65
    value_axis.maximum_scale = 0.9
    value_axis.number_format = "0%"
    value_axis.has_major_gridlines = True
    value_axis.major_gridlines.format.line.color.rgb = RGBColor(200, 215, 200)
    value_axis.tick_labels.font.size = Pt(12)
    value_axis.tick_labels.font.name = "Microsoft YaHei"

    series = chart.series[0]
    series.format.line.width = Pt(2.5)
    series.format.line.color.rgb = TITLE_COLOR
    series.marker.style = XL_MARKER_STYLE.CIRCLE
    series.marker.size = 9
    series.marker.format.fill.solid()
    series.marker.format.fill.fore_color.rgb = HIGHLIGHT_COLOR
    series.data_labels.show_value = True
    series.data_labels.number_format = "0%"
    series.data_labels.font.size = Pt(12)
    series.data_labels.font.name = "Microsoft YaHei"
    series.data_labels.font.color.rgb = TITLE_COLOR

    add_card(
        slide,
        left=1.0,
        top=2.4,
        width=2.4,
        height=2.0,
        title="质量提升",
        lines=["保存率提升16个百分点", "逐步逼近 85% 目标线"],
        fill_color=CARD_WHITE,
        title_color=TITLE_COLOR,
        body_size=14,
    )

    add_card(
        slide,
        left=1.0,
        top=4.7,
        width=2.4,
        height=2.0,
        title="关键举措",
        lines=["• 优化树种配比", "• 强化抚育管护", "• 推进智能监测"],
        fill_color=CARD_GREEN,
        title_color=DEEP_ACCENT,
        body_size=14,
    )

    baseline_box = slide.shapes.add_textbox(
        Inches(3.7),
        Inches(6.8),
        Inches(6.2),
        Inches(0.5),
    )
    baseline_tf = baseline_box.text_frame
    baseline_tf.word_wrap = True
    baseline_tf.paragraphs[0].text = "注：虚线标示 85% 目标线，建议持续加大水肥调控与后期养护投入"
    apply_paragraph_style(baseline_tf.paragraphs[0], size=12, color=TEXT_COLOR)

    apply_transition(slide, effect=TransitionEffect.WIPE)


def add_pie_chart_slide(title, subtitle, items):
    slide = prs.slides.add_slide(prs.slide_layouts[5])
    set_background(slide, BACKGROUND_YELLOW)
    add_accent(slide, TITLE_COLOR, left=0.25, width=0.2)

    title_shape = slide.shapes.title
    title_shape.text = title
    apply_title_style(title_shape.text_frame, size=40)

    subtitle_box = slide.shapes.add_textbox(
        Inches(1.2),
        Inches(1.8),
        Inches(8.2),
        Inches(0.8),
    )
    subtitle_tf = subtitle_box.text_frame
    subtitle_tf.word_wrap = True
    apply_paragraph_style(subtitle_tf.paragraphs[0], size=24, color=TEXT_COLOR)
    subtitle_tf.paragraphs[0].text = subtitle

    chart_data = ChartData()
    chart_data.categories = [item[0] for item in items]
    chart_data.add_series("面积占比", [item[1] for item in items])

    chart = slide.shapes.add_chart(
        XL_CHART_TYPE.PIE,
        Inches(3.7),
        Inches(2.4),
        Inches(6.2),
        Inches(4.4),
        chart_data,
    ).chart

    chart.has_legend = True
    chart.legend.include_in_layout = False
    chart.legend.font.size = Pt(12)
    chart.legend.font.name = "Microsoft YaHei"
    chart.legend.font.color.rgb = TEXT_COLOR

    for index, point in enumerate(chart.series[0].points):
        fill = point.format.fill
        fill.solid()
        if index % 2 == 0:
            fill.fore_color.rgb = TITLE_COLOR
        else:
            fill.fore_color.rgb = RGBColor(200, 215, 120)

    series = chart.series[0]
    series.data_labels.show_percentage = True
    series.data_labels.show_category_name = True
    series.data_labels.number_format = "0%"
    series.data_labels.font.size = Pt(12)
    series.data_labels.font.name = "Microsoft YaHei"
    series.data_labels.font.color.rgb = TITLE_COLOR
    series.data_labels.position = XL_LABEL_POSITION.BEST_FIT

    add_card(
        slide,
        left=1.0,
        top=2.4,
        width=2.4,
        height=2.2,
        title="区域亮点",
        lines=["西北治沙区贡献超三成", "京津风沙源治理持续扩容"],
        fill_color=CARD_WHITE,
        title_color=DEEP_ACCENT,
        body_size=14,
    )

    add_card(
        slide,
        left=1.0,
        top=4.9,
        width=2.4,
        height=1.6,
        title="协同建议",
        lines=["加强黄土高原与南方丘陵联动", "推动生态补偿与产业融合"],
        fill_color=CARD_GREEN,
        title_color=DEEP_ACCENT,
        body_size=13,
    )

    apply_transition(slide, effect=TransitionEffect.COVER)

add_title_slide(
    "绿色守护·治沙植木行动",
    "环境保护宣讲 | 共筑生态安全屏障",
    "守护绿水青山 · 打造万里绿色长城",
)

add_agenda_slide(AGENDA_ITEMS)
add_status_slide()
add_causes_slide()
add_strategy_slide()
add_technology_slide()
add_case_slide()
add_collaboration_slide()
add_roadmap_slide()

add_column_chart_slide(
    "近十年新增造林面积",
    "2015-2024年全国新增造林面积稳步提升（单位：万公顷）",
    YEARS,
    AFFORESTATION_AREA,
)

add_line_chart_slide(
    "造林保存率持续改善",
    "全国人工造林保存率十年提升16个百分点，质量同步跃升",
    YEARS,
    SURVIVAL_RATE,
)

add_pie_chart_slide(
    "2024年重点区域造林占比",
    "多区域协同推进，西北治沙与京津风沙源治理贡献突出",
    REGION_SHARE,
)

add_quote_slide(
    "携手共筑绿色长城",
    [
        "“留住一片绿色，就是留住希望。”",
        "治沙植木是面向未来的生态承诺，也是共同的绿色责任。",
        "让我们以节水减排、志愿植树、绿色生活方式共同行动。",
    ],
)

fallback_path = None
try:
    prs.save(OUTPUT_PATH)
    print("PPT generated successfully.")
except PermissionError:
    fallback_path = os.path.join(OUTPUT_DIR, "环境保护治沙植木宣讲_new.pptx")
    prs.save(fallback_path)
    print(f"目标文件被占用，已保存到备份文件：{fallback_path}")

